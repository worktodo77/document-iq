"""Per-page text extraction — vendored from MIP 3.9 ``api/docs_extract.py``.

The §11 reuse audit ruled REUSE on that module, so its hard-won behavior is
kept verbatim: hybrid per-page native/OCR routing, content-sniffing recovery
for misnamed production files, ZIP anti-DoS guards, per-page OCR failure
isolation, XLSX/CSV row caps with disclosed truncation, CSV encoding and
delimiter fallback, and the deterministic ``[PHOTO]`` EXIF block.

What changed, and why each change was mandatory rather than cosmetic:

* **Returns pages, not one joined string.** The original joined pages with an
  inline ``[page N]`` marker. DocIQ needs per-page OCR confidence (§4 Stage 2)
  and per-page KEEP/DROP (§4 Stage 4), and the freeze forbids a marker inside
  ``PageRecord.text`` — markers are rendered by ``emit/cleantext.py`` alone.
* **No network.** The original's ``_ocr_engine`` called ``enable_os_trust()``
  so a one-time OCR model download could get through a corporate proxy.
  Principle 4 admits no network call at all. Models are loaded from the
  installed ``rapidocr_onnxruntime`` package directory and their absence is a
  loud, actionable failure — never a download.
* **Per-line OCR confidence is captured.** The original discarded ``line[2]``,
  which is exactly the number §4 Stage 2 needs.
* **No extract cache.** The original persisted extracted text under
  ``~/.mip39/`` — outside the working folder, which §10 forbids. See
  :ref:`no-cache` below for why it was removed rather than relocated.
* **No AI captioning.** The original could call a local vision model to
  describe a photo. §12 puts any AI processing out of scope and the contract's
  ``PageKind.PHOTO`` says "never AI-captioned in DocIQ".
* **Scratch files stay inside the working folder.** ``.msg`` parsing needs a
  real file on disk; the caller supplies where.

.. _no-cache:

**Why the cache was removed, not relocated.** Relocating it under the output
root would satisfy §10's letter and defeat the determinism proof: runs 2..N of
an identical-input repeat would replay cached bytes instead of re-extracting,
so the proof would demonstrate that a cache is a cache. A content-hash-keyed
cache is also a standing correctness hazard — a stale entry written by a
different engine version replays old text under a new run's identity. The
expensive path (OCR) is the one worth caching, and Sprint 2 can reintroduce
caching inside the matter folder *behind* a flag that the determinism harness
disables. Nothing depends on it today.
"""

from __future__ import annotations

import datetime
import hashlib
import io
import os
import re
import sys
import threading
from dataclasses import dataclass
from pathlib import Path

from ..contracts import ExtractionError, PageKind, PageRecord, ProcessingStatus
from ..identify.bates import FOOTER_BLOCK_MAX_LINES
from .pagemodel import M_OCR_BLANK, make_page, synthetic_pages

# Tier 1 (§3) — extracted page by page.
TIER1_EXTENSIONS = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".xlsx": "Excel",
    ".xlsm": "Excel",
    ".xls": "Excel (legacy)",
    ".csv": "CSV",
    ".txt": "Text",
    ".md": "Markdown",
    ".log": "Log",
    ".eml": "Email",
    ".email": "Email",
    ".msg": "Outlook email",
    ".pptx": "Presentation",
    ".zip": "Archive",
}

# Tier 2 (§3, as amended by D-02 and D-10) — inventoried and hashed, never
# extracted, never blocking. The remediation hint is what turns "unsupported"
# into an action the operator can take.
TIER2_EXTENSIONS = {
    ".doc": "Legacy Word — open in Word and Save-As .docx or .pdf to include",
    ".rtf": "Rich Text — open in Word and Save-As .docx or .pdf to include",
    ".xer": "Primavera P6 export — schedule data, out of scope for v1",
    ".mpp": "Microsoft Project — export to PDF to include",
    ".dwg": "CAD drawing — plot to PDF to include",
    ".rar": "RAR archive — re-pack as .zip to include its members",
    ".png": "Image — no text layer; §3 lists images, they are not OCR'd as documents",
    ".jpg": "Image — no text layer; §3 lists images, they are not OCR'd as documents",
    ".jpeg": "Image — no text layer; §3 lists images, they are not OCR'd as documents",
    ".tif": "Image — no text layer; §3 lists images, they are not OCR'd as documents",
    ".tiff": "Image — no text layer; §3 lists images, they are not OCR'd as documents",
    ".bmp": "Image — no text layer; §3 lists images, they are not OCR'd as documents",
}

UNKNOWN_HINT = "Unrecognized format — inventoried and hashed only"

# A page with less native text than this is treated as image-only and routed to
# OCR. Inherited verbatim from the MIP 3.9 extractor, where it was measured
# against real scanned productions: a genuine text page clears it trivially,
# and a scanned page's stray header text does not.
_NATIVE_TEXT_FLOOR = 40

# ---------------------------------------------------------------------------
# Degradation markers — the one list of "this document did not read cleanly"
# ---------------------------------------------------------------------------
#
# Every place in this module that swallows an exception and carries on with
# LESS content than the source holds emits a note containing exactly one of
# these phrases. The walker's serial-retry pass keys off them, so the phrases
# are constants used in the f-strings rather than prose repeated by hand: a new
# degradation path that forgets to use one is invisible to the retry, and a
# reworded note that drifts from the matcher is the same defect a year later.
#
# A whole-document FAILED status is deliberately NOT in this list — it is a
# status, not a note, and the walker tests it directly.

M_OCR_PAGE = "ocr: page failed"
M_OCR_DOC = "OCR pass failed for this document"
M_OCR_FOOTER = "footer re-OCR pass failed"
M_ATTACH_ENUM = "could not enumerate attachments"
M_MSG_ATTACH = "could not read .msg attachments"
M_ATTACH_READ = "an attachment could not be read"
M_ZIP_MEMBER = "archive member unreadable"
M_ZIP_ATTACH = "attached archive unreadable"
M_PHOTO_PROBE = "image/EXIF probe failed"
M_SLIDE_NOTES = "slide notes could not be read"
M_EML_BODY = "email body could not be read"

TRANSIENT_MARKERS: tuple[str, ...] = (
    M_OCR_PAGE,
    M_OCR_DOC,
    M_OCR_FOOTER,
    M_ATTACH_ENUM,
    M_MSG_ATTACH,
    M_ATTACH_READ,
    M_ZIP_MEMBER,
    M_ZIP_ATTACH,
    M_PHOTO_PROBE,
    M_SLIDE_NOTES,
    M_EML_BODY,
)

# The second half of the vocabulary (Codex review #1, B-3).
#
# Some evidence gaps are not worth re-reading for: the same bytes re-parsed the
# same way will reach the same wall. An embedded .msg that this pass cannot
# flatten, a MIME part with no decodable payload, an email whose envelope will
# not parse at all — none of those are races, and putting them in
# TRANSIENT_MARKERS would spend the whole retry budget proving it.
#
# They are still evidence gaps, and B-3's requirement is that EVERY exception
# path yielding less evidence carries a marker, not only the retryable ones.
# Without a final vocabulary the choice was between "retry something that
# cannot improve" and "disclose in prose that nothing downstream can find
# mechanically", and the second is how these paths went unnoticed.

M_EML_PARSE = "email envelope would not parse"
M_ATTACH_SKIPPED = "attachment content was not brought in"

FINAL_MARKERS: tuple[str, ...] = (
    M_EML_PARSE,
    M_ATTACH_SKIPPED,
    # Routed to OCR, engine ran, nothing came back. Defined in
    # :mod:`dociq.ingest.pagemodel` and imported here so it is classified with
    # the rest of the vocabulary rather than living outside it — the class
    # assertion in the extraction tests is what surfaced the omission.
    #
    # FINAL, not transient: the same bytes through the same engine reach the
    # same wall, and a corpus of blank scans would otherwise spend the whole
    # serial-retry budget proving it. It is the one marker whose meaning is
    # ambiguous per page and unambiguous in bulk — one page is an unreadable
    # scan, every page is a dead engine — which is exactly what
    # :func:`ocr_yield` is for.
    M_OCR_BLANK,
)


def has_transient_marker(text: str | None) -> bool:
    """True when a note says this document read with less than it holds, and a
    re-read alone might recover it.

    "Transient" is the possibility being tested, not a claim: the same phrase
    covers a permanently corrupt page and a page that lost a race under load,
    and telling them apart is exactly what the walker's serial retry does.
    """
    return bool(text) and any(m in text for m in TRANSIENT_MARKERS)


def has_final_marker(text: str | None) -> bool:
    """True when a note says evidence is missing and a re-read will not help.

    Not retried, still audited: :mod:`dociq.verify.accounting` counts these so
    a final gap is a number on the run's own accounting line rather than one
    sentence inside one document's notes.
    """
    return bool(text) and any(m in text for m in FINAL_MARKERS)


def has_evidence_marker(text: str | None) -> bool:
    """True when a note says ANY evidence is missing, transient or final.

    The check a caller wants when the question is "did this document read
    completely", as opposed to "should this document be re-read".
    """
    return has_transient_marker(text) or has_final_marker(text)


_XLSX_MAX_ROWS = int(os.environ.get("DOCIQ_XLSX_MAX_ROWS", "50000"))
_CSV_MAX_ROWS = int(os.environ.get("DOCIQ_CSV_MAX_ROWS", "50000"))
_ZIP_MAX_MB = int(os.environ.get("DOCIQ_ZIP_MAX_MB", "500"))
_ZIP_MAX_MEMBERS = int(os.environ.get("DOCIQ_ZIP_MAX_MEMBERS", "2000"))
_ZIP_MAX_DEPTH = int(os.environ.get("DOCIQ_ZIP_MAX_DEPTH", "3"))


def effective_caps() -> dict[str, int]:
    """The caps this process will actually apply, for the run identity.

    Read from the same module-level constants the extractors use, not from the
    environment a second time: a second read could disagree with the first if
    the environment changed after import, and the identity must record what the
    run *did*, not what it was asked to do.

    Codex review #1 finding B-2: when one of these bites, the same folder,
    profile and index produce different evidence under an identical hashed
    configuration. Per-document truncation notes disclose the effect; they do
    not repair the identity.
    """
    return {
        "xlsx_max_rows": _XLSX_MAX_ROWS,
        "csv_max_rows": _CSV_MAX_ROWS,
        "zip_max_mb": _ZIP_MAX_MB,
        "zip_max_members": _ZIP_MAX_MEMBERS,
        "zip_max_depth": _ZIP_MAX_DEPTH,
    }


@dataclass(frozen=True, slots=True)
class ExtractedDoc:
    """What one Tier-1 file yielded. Never raises out of :func:`extract`."""

    pages: tuple[PageRecord, ...] = ()
    notes: tuple[str, ...] = ()
    status: ProcessingStatus = ProcessingStatus.FULL
    error: str | None = None


@dataclass
class ExtractOptions:
    """Everything that can change extracted bytes, passed explicitly.

    Module-level mutable state would make two concurrently-running matters
    influence each other's output, which the determinism contract forbids.
    """

    conf_threshold: float = 0.85
    scratch_dir: Path | None = None
    """Where formats that need a real file on disk may write one. §10 forbids
    persistent temp files outside the working folder, so the caller — which
    knows the output root — decides. ``None`` falls back to the system temp
    directory and the file is unlinked in a ``finally``."""

    ocr_enabled: bool = True
    """Off only for tests that must exercise the native path in isolation."""

    footer_reocr: bool = True
    """The D-25 targeted footer re-OCR.

    A field rather than a module flag for the same reason as everything else
    here: two matters running concurrently must not be able to change each
    other's bytes. Off only for the tests that measure what the pass costs and
    what it recovers, which need the before as well as the after."""


# ---------------------------------------------------------------------------
# OCR — engine, models, and per-line confidence
# ---------------------------------------------------------------------------

_OCR_ENGINE = None
_OCR_LOCK = threading.Lock()
_OCR_POOL = None

_OCR_PAGE_WORKERS = int(os.environ.get(
    "DOCIQ_OCR_WORKERS", str(min(16, max(1, (os.cpu_count() or 2) - 2)))))

_MODEL_FILES = (
    "ch_PP-OCRv3_det_infer.onnx",
    "ch_PP-OCRv3_rec_infer.onnx",
    "ch_ppocr_mobile_v2.0_cls_infer.onnx",
)


def ocr_model_dir() -> Path:
    """Directory holding the bundled ONNX models.

    ``DOCIQ_OCR_MODEL_DIR`` overrides it, which is how an operator points a run
    at models held somewhere else. The default is the installed package's own
    ``models/`` directory — the wheel ships them, so a correct install already
    has every byte the OCR path needs and nothing is ever fetched.

    **Frozen builds resolve it from the bundle, not from ``__file__``.** In a
    PyInstaller build ``rapidocr_onnxruntime.__file__`` names a path inside the
    PYZ archive that does not exist on disk; the models are real files unpacked
    under ``sys._MEIPASS``. The frozen branch is explicit rather than relying on
    the two happening to coincide, because if they ever stop coinciding the
    failure is "OCR unavailable" on a client machine and nowhere else.
    """
    override = os.environ.get("DOCIQ_OCR_MODEL_DIR")
    if override:
        return Path(override)
    meipass = getattr(sys, "_MEIPASS", None)
    if meipass:
        return Path(meipass) / "rapidocr_onnxruntime" / "models"
    import rapidocr_onnxruntime

    return Path(rapidocr_onnxruntime.__file__).parent / "models"


def ocr_models_present() -> tuple[bool, str]:
    """``(ok, message)``. The message names the missing file and the fix."""
    try:
        d = ocr_model_dir()
    except Exception as exc:
        return False, f"OCR engine 'rapidocr_onnxruntime' is not importable: {exc}"
    missing = [m for m in _MODEL_FILES if not (d / m).is_file()]
    if missing:
        return False, (
            f"OCR models are missing from {d}: {', '.join(missing)}. "
            "DocIQ never downloads them (Principle 4 — no network). Reinstall "
            "rapidocr-onnxruntime, or set DOCIQ_OCR_MODEL_DIR to a directory "
            "containing the three .onnx model files."
        )
    return True, ""


def ocr_available() -> bool:
    """True when the local OCR stack is importable AND its models are on disk.

    **A PRESENCE check, not a capability check — the name overstates it and the
    docstring is where that stops.** It imports two modules and stats three
    ``.onnx`` files. It never constructs the engine and never runs inference,
    so an engine that imports cleanly, finds its models and then produces
    nothing on every page passes it. That is where the Sprint-2 burn happened:
    *inside* inference, under :func:`_ocr_pdf_pages`'s per-page
    ``except Exception``, where a totally dead engine and a few bad pages look
    identical page by page.

    The real capability check is :func:`dociq.selftest`'s cold-construction
    probe, which builds the engine and OCRs a synthetic image — and
    ``build.py --skip-verify`` bypasses it. The run-level backstop for a build
    that shipped anyway is :func:`ocr_yield`, whose whole subject is the case
    this function cannot see.
    """
    try:
        import fitz  # noqa: F401  (pymupdf)
        import rapidocr_onnxruntime  # noqa: F401
    except Exception:
        return False
    return ocr_models_present()[0]


OCR_DEAD_ENGINE = (
    "OCR produced no text on ANY of the {attempted:,} page(s) it was run on in "
    "this run. A single page that recovers nothing is ordinary — a blank or "
    "unreadable scan — but every page recovering nothing is what a dead OCR "
    "engine looks like from the outside, and the per-page notes offer the "
    "innocent explanation first. Before relying on this corpus, run "
    "`dociq selftest` (it builds the engine and OCRs a test image); if that "
    "passes, these pages really are unreadable and the run stands."
)
"""The run-level alarm §4's per-page notes structurally cannot raise.

Per-document notes say "N page(s) routed to OCR recovered no text" and are read
one document at a time, where "some bad scans" is the natural reading and is
usually right. Nothing was looking at the whole run, which is the only scale at
which "every attempt, without exception" is visible — and that is the shape of
a dead engine rather than of bad pages.
"""


def ocr_yield(documents) -> tuple[int, int]:
    """``(pages OCR was attempted on, pages that recovered text)`` for a corpus.

    Reconstructed from the final page records rather than from a counter, so it
    describes the deliverable — including after a serial retry replaced a
    document's records wholesale.

    A page counts as ATTEMPTED on the DISCLOSURE, not on the kind. A page routed
    to OCR that recovers nothing is re-labelled ``EMPTY`` by
    :func:`dociq.ingest.pagemodel.make_page` (``EMPTY`` is the only kind the
    contract lets carry no ``ocr_conf``), and page 1 of a photo-only document is
    re-labelled ``PHOTO`` before that. Counting kinds would therefore have
    counted zero attempts on precisely the run this exists to catch — measured,
    not reasoned: a dead-engine walk over the scanned fixture yielded one PHOTO
    page and one EMPTY page and no ``PageKind.OCR`` at all. The three
    disclosures below survive both relabellings, which is why they are the
    thing counted:

    * :data:`~dociq.ingest.pagemodel.M_OCR_BLANK` — routed to OCR, recovered
      nothing;
    * :data:`M_OCR_PAGE` — could not even be rasterized or read;
    * ``PageKind.OCR`` — recovered text, i.e. the attempts that worked.
    """
    attempted = recovered = 0
    for doc in documents:
        for page in doc.pages:
            worked = page.kind is PageKind.OCR and page.text.strip()
            blank = any(n.startswith(M_OCR_BLANK) or n.startswith(M_OCR_PAGE)
                        for n in page.notes)
            if not (worked or blank):
                continue
            attempted += 1
            if worked:
                recovered += 1
    return attempted, recovered


def ocr_yield_warning(documents) -> str | None:
    """:data:`OCR_DEAD_ENGINE`, filled in, when a run recovered nothing at all."""
    attempted, recovered = ocr_yield(documents)
    if attempted and not recovered:
        return OCR_DEAD_ENGINE.format(attempted=attempted)
    return None


def ocr_engine_version() -> str:
    """Recorded in ``RunConfig`` — the engine identity is part of run identity."""
    try:
        from importlib.metadata import version

        return version("rapidocr_onnxruntime")
    except Exception:
        return "unknown"


_MODEL_ID_CACHE: dict[tuple, str] = {}
_MODEL_ID_LOCK = threading.Lock()


def _model_stat_key(d: Path) -> tuple:
    """Cheap identity of the model directory: name, size and mtime per file.

    Used only as a *cache key*, never as the recorded identity. Sizes and mtimes
    are what a stale cache would miss; the recorded identity is always the
    content hash below.
    """
    out = []
    for name in _MODEL_FILES:
        p = d / name
        try:
            st = p.stat()
            out.append((name, st.st_size, st.st_mtime_ns))
        except OSError:
            out.append((name, -1, -1))
    return (str(d),) + tuple(out)


def ocr_model_id() -> str:
    """Stable identity of the OCR model artifacts — package version PLUS a hash
    of the model bytes.

    Recorded in ``RunConfig.limits.ocr_model_id`` and therefore in the hashed
    run identity (Codex review #1 finding B-2). A version string alone is not
    enough: ``DOCIQ_OCR_MODEL_DIR`` can point the same installed package at
    different ONNX files, and two engines that read a page differently are
    different inputs to the run, however they were installed.

    The hash is over the three model files' names and bytes in a fixed order, so
    it is independent of directory listing order and of where the files live.
    Nothing is downloaded — Principle 4 — and a missing or unreadable model
    yields an explicit ``models-unavailable`` identity rather than a silent
    empty string that would compare equal to a run that had no OCR at all.
    """
    version = ocr_engine_version()
    try:
        d = ocr_model_dir()
    except Exception:
        return f"rapidocr_onnxruntime {version}; models-unavailable"
    key = _model_stat_key(d)
    with _MODEL_ID_LOCK:
        hit = _MODEL_ID_CACHE.get(key)
    if hit is not None:
        return hit
    h = hashlib.sha256()
    try:
        for name in _MODEL_FILES:
            h.update(name.encode("utf-8"))
            h.update(b"\0")
            with open(d / name, "rb") as fh:
                while chunk := fh.read(1 << 20):
                    h.update(chunk)
        ident = f"rapidocr_onnxruntime {version}; models {h.hexdigest()[:32]}"
    except OSError:
        ident = f"rapidocr_onnxruntime {version}; models-unavailable"
    with _MODEL_ID_LOCK:
        _MODEL_ID_CACHE[key] = ident
    return ident


def _ocr_engine():
    """The shared RapidOCR engine.

    Construction is LOCKED, inherited from MIP 3.9: the page pool can hit this
    from ~16 threads at once and an unlocked lazy init builds N redundant
    engines (~200 MB of ONNX each — 1.5 GB RSS observed on a 7-PDF batch).

    Model paths are passed explicitly rather than left to the library's
    relative-path default, so a caller running from another working directory
    cannot silently miss them and trigger a fetch attempt.
    """
    global _OCR_ENGINE
    if _OCR_ENGINE is not None:
        return _OCR_ENGINE
    with _OCR_LOCK:
        if _OCR_ENGINE is None:
            ok, msg = ocr_models_present()
            if not ok:
                raise ExtractionError(msg)
            from rapidocr_onnxruntime import RapidOCR

            d = ocr_model_dir()
            _OCR_ENGINE = RapidOCR(
                det_model_path=str(d / _MODEL_FILES[0]),
                rec_model_path=str(d / _MODEL_FILES[1]),
                cls_model_path=str(d / _MODEL_FILES[2]),
            )
    return _OCR_ENGINE


@dataclass(frozen=True, slots=True)
class OcrLine:
    """One recognized line: where it was, what it said, how sure the engine was.

    The box is carried even though no contract field holds it, because §4
    Stage 3 matches Bates stamps against page *corners and footers* — that is a
    geometry question, and Track B would otherwise have to stand up a second
    OCR engine to ask it. Nothing here reaches disk.
    """

    box: tuple[tuple[float, float], ...]
    text: str
    conf: float


def ocr_lines(arr) -> list[OcrLine]:
    """OCR one image array into per-line results.

    The confidence is the whole point of the change from MIP 3.9, which joined
    ``line[1]`` and dropped ``line[2]`` — and ``line[2]`` is exactly what §4
    Stage 2's 85% review threshold is measured against.
    """
    res, _ = _ocr_engine()(arr)
    out: list[OcrLine] = []
    for line in res or []:
        try:
            conf = float(line[2])
        except (IndexError, TypeError, ValueError):
            # A line without a usable score is not silently perfect: score it
            # zero so it counts against the mean and lands in the low-confidence
            # tally, which is the direction that gets a human to look.
            conf = 0.0
        try:
            box = tuple((float(x), float(y)) for x, y in line[0])
        except (IndexError, TypeError, ValueError):
            box = ()
        out.append(OcrLine(box=box, text=str(line[1]), conf=conf))
    return out


def _ocr_array(arr) -> tuple[str, list[float]]:
    """``(joined text, per-line confidences)`` — what the page model needs."""
    lines = ocr_lines(arr)
    return " ".join(ln.text for ln in lines), [ln.conf for ln in lines]


@dataclass(frozen=True, slots=True)
class _OcrPage:
    text: str = ""
    confs: tuple[float, ...] = ()
    failed: bool = False


def _page_array(page):
    """Rasterize one PDF page straight into an OCR-ready BGR array.

    No PNG round-trip — the encode/decode costs ~0.2 s/page for nothing.
    """
    import cv2
    import numpy as np

    pix = page.get_pixmap(dpi=200)
    arr = np.frombuffer(pix.samples, np.uint8).reshape(pix.height, pix.width, pix.n)
    if pix.n == 1:
        return cv2.cvtColor(arr, cv2.COLOR_GRAY2BGR)
    if pix.n == 3:
        return cv2.cvtColor(arr, cv2.COLOR_RGB2BGR)
    return cv2.cvtColor(arr, cv2.COLOR_RGBA2BGR)


def _ocr_page_pool():
    global _OCR_POOL
    if _OCR_POOL is None:
        with _OCR_LOCK:
            if _OCR_POOL is None:
                from concurrent.futures import ThreadPoolExecutor

                _OCR_POOL = ThreadPoolExecutor(
                    max_workers=_OCR_PAGE_WORKERS, thread_name_prefix="ocr-page")
    return _OCR_POOL


def _ocr_pdf_pages(raw: bytes, pages: list[int]) -> dict[int, _OcrPage]:
    """OCR only the given 0-based page indices, fanned across the shared pool.

    Results are keyed by page index and the caller reassembles by index, never
    by completion order — the pool returns pages as they finish, and an
    order-of-completion assembly would produce a different document on every
    run. The contract's gapless-1..N check would catch a *missing* page but not
    a *permuted* one, so this is the primary defence, not the backstop.

    Rasterization is chunked: a 300-page scan at 200 dpi is ~11 MB/page, so
    rasterizing everything up front would cost gigabytes per file.
    """
    import fitz  # pymupdf

    out: dict[int, _OcrPage] = {}
    chunk_n = 16
    pool = _ocr_page_pool()
    with fitz.open(stream=raw, filetype="pdf") as doc:
        idxs = [i for i in pages if 0 <= i < len(doc)]
        for c0 in range(0, len(idxs), chunk_n):
            arrays: dict[int, object] = {}
            for i in idxs[c0:c0 + chunk_n]:
                try:
                    arrays[i] = _page_array(doc[i])
                except Exception:
                    out[i] = _OcrPage(failed=True)  # one bad page must not sink the doc
            futs = {i: pool.submit(_ocr_array, a) for i, a in arrays.items()}
            for i, fut in futs.items():
                try:
                    text, confs = fut.result()
                    out[i] = _OcrPage(text=text, confs=tuple(confs))
                except Exception:
                    out[i] = _OcrPage(failed=True)
    return out


# ---------------------------------------------------------------------------
# Targeted footer re-OCR — D-25
# ---------------------------------------------------------------------------
#
# A Bates stamp is not body text and reading it as body text is what the
# criterion-4 acceptance run measured the cost of: 100.000% of native-text pages
# carried their locator and 31.250% of OCR'd pages did (593/648 overall against
# a >=99% bar, 0 wrong, 0 false positives, 55 absent). The whole-page pass is
# one recognition tuned for a page of prose; a six-character stamp in a 10pt
# footer is a rounding error inside it.
#
# So the stamp gets its own recognition. The band of the page where a stamp
# lives is rasterized on its own at a much higher resolution and read again,
# and only the stamp-shaped tokens of that reading are merged back.
#
# THREE PROPERTIES ARE LOAD-BEARING AND EACH IS ENFORCED HERE RATHER THAN
# HOPED FOR:
#
# 1. It runs ONLY where it can help. The pass fires on a page that (a) DocIQ
#    had to OCR at all and (b) whose ordinary reading yielded no stamp-shaped
#    line anywhere in the Bates zone. A native-text page never pays for it, and
#    on the acceptance corpus that is 568 of 648 pages.
# 2. It cannot turn a miss into a WRONG answer. Nothing here writes a locator.
#    It appends candidate TEXT, which Stage 3 then judges with exactly the same
#    grammar and the same operator-confirmed format as any other text. A
#    misread footer produces a token that does not match the confirmed format
#    and is ignored, which leaves the page a flagged miss.
# 3. Nothing about the attempt reaches hashed content. No timing, no retry
#    count, no resolution, no marker. The output is a deterministic function of
#    the page bytes: the same PDF yields the same appended tokens forever.

FOOTER_REOCR_DPI = 400
"""Rasterization resolution for the band pass, against 200 dpi for the page.

Doubling the resolution doubles the glyph height the recognizer sees BEFORE its
own fixed-height crop resize, which is the mechanism: at 200 dpi a 10pt stamp is
~28 px tall and the recognizer's input is 48 px, so it is upsampling guesswork;
at 400 dpi it is downsampling a real reading.

**Measured, and the measurement bounds the claim.** Over the 55 pages the
criterion-4 baseline missed, 400 dpi reads the stamp's DIGITS correctly where
the whole-page pass read nothing at all. 300 dpi is close behind; 600 and 800
dpi are WORSE, not better, which is the opposite of the naive expectation and
the reason this number is measured rather than maximised. See
``docs/verification/bates_d25_2026-08-01.md`` for the sweep, and for the part of
the stamp the band pass does *not* fix.
"""

FOOTER_REOCR_BAND_PT = 90.0
"""Height of the stamp band, in PDF points — a physical 1.25 inches, not a
fraction of the page.

This was a fraction of page height in the first draft and that was a defect,
found by measuring rather than by reading. The acceptance corpus contains pages
that are **2700 x 3681 points** — photographs whose page box is 37 x 51 inches —
and 14% of that page is a seven-inch strip. Two things went wrong at once: the
band became tens of megapixels (a sweep over 55 such pages did not finish in 90
minutes), and its aspect ratio tripped the bypass below.

A Bates stamp is burned in at a physical size, a physical distance from the
physical edge of the page. On the acceptance corpus it sits within ~40 points of
the bottom edge whether the page is letter-size or four feet tall. So the band
is physical too, and its cost stops scaling with page AREA.
"""

_FOOTER_MAX_ASPECT = 8.0
"""rapidocr SKIPS text detection entirely when width/height exceeds its
``width_height_ratio`` (8.0 in the shipped ``config.yaml``) and recognizes the
whole strip as a single line. Measured on the corpus: a 3600pt-wide page's
footer strip is 13:1, detection was bypassed, and the pass returned nothing
readable at all. So the band is TILED — never stretched — and every tile is at
most this ratio wide. It is the engine's own configuration, not a taste."""

_FOOTER_TILE_OVERLAP = 0.2
"""How much consecutive tiles overlap, as a fraction of tile width.

Without it a stamp straddling a tile boundary is cut in half and read as two
things, neither of which is a locator. One fifth is comfortably wider than any
stamp relative to a ten-inch tile."""

_FOOTER_MAX_TILES = 12
"""Hard ceiling on tiles per band, disclosed rather than silent.

Tile width is ``8 x band`` and ``band`` is ``min(1.25in, page height)``, so a
page that is very wide and very SHORT makes tiles arbitrarily narrow: a
612 x 0.1pt page yields nearly a thousand of them. That is a degenerate page
rather than a real one, but "degenerate input cannot reach this" is exactly the
assumption that produces a run that never finishes.

Twelve covers a page ten feet wide at the full band height, which is past any
sheet size a production contains. Beyond it the band is read left to right and
the REMAINDER IS NOT READ — a stamp out there is a miss, which is the failure
direction §4 asks for, and the ceiling is stated here rather than discovered."""

_FOOTER_PROBE_PAGES = 4
"""How many of a document's qualifying pages are probed on BOTH bands before
the pass decides which band, if either, reads this production.

Four rather than one because a production's first scanned page can be a cover
sheet, a photograph, or a fax header, and one unlucky page should not switch the
pass off for a 300-page document. Four rather than forty because the decision is
about where the production burns its stamp, and that does not vary page to page.

Stated rather than tuned: a stamped production resolves the very first page
probed, so this bound costs it nothing; an unreadable scan pays eight
recognitions instead of thousands, and the pages it declines to read are
reported in the document's notes."""

_FOOTER_CHUNK_PAGES = 8
"""Pages whose tiles are rasterized before any of them is recognized. Bounds
peak memory the way :func:`_ocr_pdf_pages` does, and lower here because one page
can be several tiles."""


def _band_tiles(page, dpi: int, band_pt: float, *, top: bool) -> list:
    """The stamp band of one page, as OCR-ready tiles in left-to-right order.

    The clip is taken in PDF user space and rendered at ``dpi``, so each tile is
    a genuine re-render at higher resolution — not an upscale of the 200 dpi
    page image, which would add no information at all.

    Cost is bounded by construction, and the bound is arithmetic rather than a
    cap: a tile is at most ``band_pt`` tall and ``_FOOTER_MAX_ASPECT * band_pt``
    wide, so a tile is ~1.7 Mpx at 400 dpi whatever the page is, and the tile
    COUNT grows with page width alone. A four-foot-wide page costs five tiles,
    not fifty megapixels.

    **The band is rendered ONCE and sliced, and that is a fix rather than a
    style.** Rendering each tile with its own ``get_pixmap`` clip re-decodes the
    page's embedded image every time: on this corpus a page is a 230 MB
    photograph, so five tiles top and bottom meant ten full decodes of it. The
    acceptance run took over an hour and a half that way and had to be killed.
    One render, then numpy views — the pixels are the same and the decode
    happens once.
    """
    import cv2
    import fitz
    import numpy as np

    r = page.rect
    band = min(band_pt, r.height)
    y0 = r.y0 if top else r.y1 - band
    pix = page.get_pixmap(dpi=dpi, clip=fitz.Rect(r.x0, y0, r.x1, y0 + band))
    arr = np.frombuffer(pix.samples, np.uint8).reshape(
        pix.height, pix.width, pix.n)
    if pix.n == 1:
        arr = cv2.cvtColor(arr, cv2.COLOR_GRAY2BGR)
    elif pix.n == 3:
        arr = cv2.cvtColor(arr, cv2.COLOR_RGB2BGR)
    else:
        arr = cv2.cvtColor(arr, cv2.COLOR_RGBA2BGR)

    h, w = arr.shape[:2]
    if h <= 0 or w <= 0:
        return []
    tile_w = min(w, max(1, int(round(_FOOTER_MAX_ASPECT * h))))
    step = max(1, int(round(tile_w * (1.0 - _FOOTER_TILE_OVERLAP))))

    spans: list[tuple[int, int]] = []
    x = 0
    while True:
        x1 = min(x + tile_w, w)
        span = (max(0, x1 - tile_w), x1)
        if span not in spans:
            spans.append(span)
        if x1 >= w or len(spans) >= _FOOTER_MAX_TILES:
            break
        x += step
    return [np.ascontiguousarray(arr[:, x0:x1]) for x0, x1 in spans]


def _band_tokens(arr) -> tuple[tuple[str, ...], tuple[float, ...]]:
    """``(stamp-shaped tokens, their confidences)`` from one band tile."""
    from ..identify.bates import stamp_tokens

    lines = ocr_lines(arr)
    tokens = stamp_tokens("\n".join(ln.text for ln in lines))
    if not tokens:
        return (), ()
    # Each token's confidence is the confidence of the recognized line it came
    # out of. A token whose line cannot be identified scores 0.0 rather than
    # silently perfect — the same rule ``ocr_lines`` applies to a line with no
    # usable score.
    confs: list[float] = []
    for tok in tokens:
        best = 0.0
        for ln in lines:
            if tok in " ".join(ln.text.split()):
                best = max(best, ln.conf)
        confs.append(best)
    return tokens, tuple(confs)


def _band_pass(doc, idxs: list[int], top: bool, pool,
               out: dict[int, tuple[tuple[str, ...], tuple[float, ...]]]) -> bool:
    """Read one band of the given pages into ``out``. True if anything read.

    **Assembled in submission order, never completion order.** Tiles are
    submitted page by page and left to right and their results are consumed in
    that same order, so which tile finishes first cannot change a page's text.
    That is the same rule, and the same reason, as :func:`_ocr_pdf_pages`, and
    it is what makes criterion 7 hold over a pass that fans out.
    """
    recovered = False
    for c0 in range(0, len(idxs), _FOOTER_CHUNK_PAGES):
        jobs: list[tuple[int, object]] = []
        for i in idxs[c0:c0 + _FOOTER_CHUNK_PAGES]:
            out.setdefault(i, ((), ()))
            try:
                tiles = _band_tiles(doc[i], FOOTER_REOCR_DPI,
                                    FOOTER_REOCR_BAND_PT, top=top)
            except Exception:
                continue          # one page that will not rasterize, not the doc
            jobs.extend((i, pool.submit(_band_tokens, a)) for a in tiles)
        merged: dict[int, tuple[list[str], list[float]]] = {}
        for i, fut in jobs:
            try:
                toks, confs = fut.result()
            except Exception:
                toks, confs = (), ()
            slot = merged.setdefault(i, ([], []))
            for tok, conf in zip(toks, confs):
                if tok in slot[0]:          # the tile overlap sees it twice
                    continue
                if len(slot[0]) >= FOOTER_BLOCK_MAX_LINES:
                    break
                slot[0].append(tok)
                slot[1].append(conf)
        for i, (toks, confs) in merged.items():
            if toks:
                out[i] = (tuple(toks), tuple(confs))
                recovered = True
    return recovered


def _reocr_bands(raw: bytes, pages: list[int]
                 ) -> dict[int, tuple[tuple[str, ...], tuple[float, ...]]]:
    """Re-read the stamp band of the given 0-based pages.

    §4 says "page corners/footers", so both the bottom and the top of the page
    are covered — a header-stamped production is an ordinary thing and assuming
    it away would be a silent limit.

    **It CALIBRATES, and that is what makes it affordable.** "Only where it can
    help" is a per-PAGE test and a per-page test is not enough: on an image-only
    document every page qualifies, and reading two bands of five tiles each on a
    300-page scan is three thousand recognitions for one document. That is not a
    hypothetical — it is what a criterion-4 re-run was doing when it passed nine
    hours and had to be killed.

    Where a production burns its stamp is a property of the PRODUCTION, not of
    the page. So the first :data:`_FOOTER_PROBE_PAGES` qualifying pages are
    probed on both bands; after that only the band(s) that actually read
    something are used, and if NEITHER read anything the pass stops for that
    document. A scan whose footers cannot be read costs eight recognitions
    instead of three thousand; a stamped production keeps every page it was
    going to get, because a production that stamps its footer resolves the very
    first page probed.

    The calibration is a function of the page contents in page order, so it is
    identical run to run for the same bytes. The pages it declines to read are
    the caller's to disclose: every page this was ASKED about has an entry in
    the returned mapping, so ``len(pages) - len(result)`` is the number skipped.
    """
    import fitz

    out: dict[int, tuple[tuple[str, ...], tuple[float, ...]]] = {}
    pool = _ocr_page_pool()
    with fitz.open(stream=raw, filetype="pdf") as doc:
        idxs = [i for i in pages if 0 <= i < len(doc)]
        if not idxs:
            return out
        probe, rest = idxs[:_FOOTER_PROBE_PAGES], idxs[_FOOTER_PROBE_PAGES:]

        # Phase 1 — probe both bands, bottom first, on the opening pages.
        productive: set[bool] = set()
        for top in (False, True):
            todo = [i for i in probe if not out.get(i, ((), ()))[0]]
            if not todo:
                break
            if _band_pass(doc, todo, top, pool, out):
                productive.add(top)

        # Phase 2 — the remainder, on the bands that demonstrably read this
        # production. Nothing productive means nothing is read at all.
        for top in (False, True):
            if top not in productive:
                continue
            todo = [i for i in rest if not out.get(i, ((), ()))[0]]
            if not todo:
                break
            _band_pass(doc, todo, top, pool, out)
    return out


# ---------------------------------------------------------------------------
# Photo PDFs — deterministic EXIF, no AI (§12)
# ---------------------------------------------------------------------------


def _gps_to_decimal(ref, vals) -> float | None:
    """One GPS coordinate as signed decimal degrees, or ``None``.

    Used to return ``0.0`` from a bare ``except``, which is not a failure value
    at all: it is a valid coordinate on the equator and on the prime meridian,
    and the caller's ``if lat or lon`` test then read an unparseable fix as "no
    fix" and dropped it with no record (Codex review #1, B-3, sibling class).
    ``None`` is unambiguous and the caller discloses it.
    """
    try:
        d, m, s = (float(v) for v in vals)
        dec = d + m / 60.0 + s / 3600.0
        return -dec if str(ref).upper() in ("S", "W") else dec
    except Exception:
        return None


def exif_from_image_bytes(img: bytes) -> tuple[dict, list[str]]:
    """``({'date': ..., 'gps': ...}, notes)`` — best-effort EXIF read.

    **Why this returns notes.** It used to return a bare dict and swallow four
    separate exceptions into ``pass``: the whole-image open, the EXIF-IFD read
    that holds ``DateTimeOriginal``, the GPS-IFD read, and each coordinate's
    own conversion. A site photo's camera date and GPS fix are the only
    evidence such a document carries — OCR never looks where a camera writes —
    so each of those silent paths deleted the entire evidentiary content of the
    page and reported a clean read. That is the same Principle-1 defect as the
    EML body walk, in the class Codex review #1 (B-3) named alongside it.

    A dict cannot say "there was EXIF here and I could not read it", so the
    return type had to widen. The notes carry :data:`M_PHOTO_PROBE`, which is
    already a transient marker, so a photo whose EXIF would not read is re-read
    serially like any other degraded document.
    """
    out: dict = {}
    notes: list[str] = []
    try:
        from PIL import ExifTags, Image

        with Image.open(io.BytesIO(img)) as im:
            exif = im.getexif()
            if not exif:
                return out, notes
            dt = None
            try:  # DateTimeOriginal lives in the EXIF IFD; fall back to DateTime
                ifd = exif.get_ifd(ExifTags.IFD.Exif)
                dt = ifd.get(ExifTags.Base.DateTimeOriginal)
            except Exception as exc:
                # Not silent even though a fallback follows: if the fallback
                # also comes back empty, the document loses its date and the
                # only reason would otherwise be invisible.
                notes.append(f"{M_PHOTO_PROBE}: the EXIF sub-directory holding "
                             f"DateTimeOriginal could not be read ({exc}); the "
                             "capture date falls back to the basic DateTime tag"
                             [:300])
            dt = dt or exif.get(ExifTags.Base.DateTime)
            if dt:
                s = str(dt).strip()  # "YYYY:MM:DD HH:MM:SS"
                if len(s) >= 10 and s[4] == ":" and s[7] == ":":
                    s = s[:10].replace(":", "-") + s[10:]
                out["date"] = s[:16]
            try:
                gps = exif.get_ifd(ExifTags.IFD.GPSInfo)
                if gps:
                    lat = _gps_to_decimal(gps.get(1), gps.get(2) or ())
                    lon = _gps_to_decimal(gps.get(3), gps.get(4) or ())
                    if lat is None or lon is None:
                        if gps.get(2) or gps.get(4):
                            notes.append(
                                f"{M_PHOTO_PROBE}: this image carries a GPS tag "
                                "whose coordinates could not be converted; the "
                                "location is absent from this document")
                    elif lat or lon:
                        out["gps"] = f"{lat:.6f}, {lon:.6f}"
            except Exception as exc:
                notes.append(f"{M_PHOTO_PROBE}: the image's GPS sub-directory "
                             f"could not be read ({exc}); the location is "
                             "absent from this document"[:300])
    except Exception as exc:
        notes.append(f"{M_PHOTO_PROBE}: the embedded image could not be opened "
                     f"for an EXIF read ({exc}); the camera date and GPS are "
                     "absent from this document"[:300])
    return out, notes


def _photo_block(raw: bytes, n_pages: int,
                 content_len: int) -> tuple[str, list[str]]:
    """``(marker text, notes)`` for an image-based PDF (a photo print-out).

    The notes are empty except when the probe itself failed. It used to return
    ``""`` from a bare ``except Exception``, which meant a file whose EXIF read
    threw — for any reason, including a transient one — silently lost its
    ``[PHOTO]`` block and its camera date, and the run said nothing at all. A
    swallowed exception that changes the emitted text is a Principle-1
    violation whatever caused it.

    Widened from one note to a list by Codex review #1 (B-3): the EXIF read
    below can now fail in four distinguishable ways and reporting only the
    first would reintroduce the same silence one level down.

    Photo test: trivial text layer plus at least one large embedded image; EXIF
    comes from the largest such image. A site photo carries its evidence where
    OCR never looks — the camera-stamped date and GPS — and surfacing it as
    text is what lets Stage 1 date the document at all.
    """
    if content_len >= max(40, 8 * n_pages):
        return "", []
    exif_notes: list[str] = []
    try:
        import fitz

        with fitz.open(stream=raw, filetype="pdf") as doc:
            biggest, n_imgs = None, 0
            for page in doc:
                for info in page.get_images(full=True):
                    n_imgs += 1
                    xref = info[0]
                    w, h = int(info[2] or 0), int(info[3] or 0)
                    if w * h < 250_000:  # ignore logos/stamps (<~0.25 MP)
                        continue
                    if biggest is None or w * h > biggest[0]:
                        biggest = (w * h, xref)
            if biggest is None:
                return "", []
            meta: dict = {}
            try:
                meta, exif_notes = exif_from_image_bytes(
                    doc.extract_image(biggest[1])["image"])
            except Exception as exc:
                exif_notes = [(f"{M_PHOTO_PROBE}: the embedded image's EXIF could "
                               f"not be read ({exc}); the camera date and GPS are "
                               "absent from this document")[:300]]
            parts = [f"[PHOTO] Image-based document ({len(doc)} page(s), "
                     f"{n_imgs} image(s))."]
            if meta.get("date"):
                parts.append(f"Camera (EXIF) date: {meta['date']}.")
            if meta.get("gps"):
                parts.append(f"GPS: {meta['gps']}.")
            parts.append("Visual content not machine-read — view the source image "
                         "for what the photo shows.")
            return " ".join(parts), exif_notes
    except Exception as exc:
        return "", [(f"{M_PHOTO_PROBE}: this document was probed as an "
                     f"image-based (photo) PDF and the probe raised ({exc}); "
                     "no [PHOTO] block was emitted")[:300]]


# ---------------------------------------------------------------------------
# Per-format extractors — each returns a list of PageRecords
# ---------------------------------------------------------------------------


def _extract_pdf(raw: bytes, opt: ExtractOptions) -> tuple[list[PageRecord], list[str]]:
    """Hybrid per-page routing: keep each page's text layer when it is
    substantive, OCR only the pages without one.

    Real productions are mixed — typed correspondence interleaved with scanned
    attachments — so this simultaneously OCRs LESS (text pages skip the
    expensive path) and extracts MORE (scanned attachments inside text-rich
    files previously contributed nothing, because the whole file took the text
    route).
    """
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover — declared, so always present
        raise ExtractionError("PDF support requires 'pypdf'.") from exc
    try:
        reader = PdfReader(io.BytesIO(raw))
        native = [(page.extract_text() or "") for page in reader.pages]
    except Exception as exc:
        raise ExtractionError(f"Could not read PDF: {exc}") from exc

    notes: list[str] = []
    n = len(native)
    # Measure actual page CONTENT so an empty text layer is not masked.
    content_len = sum(len(t.strip()) for t in native)
    photo, photo_notes = _photo_block(raw, n, content_len)
    notes.extend(photo_notes)

    ocr_by_page: dict[int, _OcrPage] = {}
    need = [i for i, t in enumerate(native) if len(t.strip()) < _NATIVE_TEXT_FLOOR]
    if need and opt.ocr_enabled:
        if not ocr_available():
            notes.append(
                f"{len(need)} page(s) have no usable text layer and OCR is "
                f"unavailable: {ocr_models_present()[1]}")
        else:
            try:
                ocr_by_page = _ocr_pdf_pages(raw, need)
            except Exception as exc:
                notes.append(f"{M_OCR_DOC}: {exc}")
    elif need and not opt.ocr_enabled:
        notes.append(f"{len(need)} page(s) have no usable text layer; OCR disabled")

    # --- D-25: the stamp gets its own recognition, where it can help --------
    # Only pages DocIQ actually OCR'd, and only those whose ordinary reading
    # produced nothing stamp-shaped anywhere in the Bates zone. The trigger is
    # a pure function of the text just read, so it is identical run to run.
    reocr: dict[int, tuple[tuple[str, ...], tuple[float, ...]]] = {}
    n_footer_declined = 0
    if ocr_by_page and opt.footer_reocr:
        from ..identify.bates import zone_has_candidate

        retry = [i for i in sorted(ocr_by_page)
                 if not ocr_by_page[i].failed
                 and not zone_has_candidate(ocr_by_page[i].text)]
        if retry:
            try:
                reocr = _reocr_bands(raw, retry)
                n_footer_declined = len(retry) - len(reocr)
            except Exception as exc:
                notes.append(f"{M_OCR_FOOTER}: {exc}")

    pages: list[PageRecord] = []
    n_ocr_failed = 0
    n_ocr_blank = 0
    n_footer_recovered = 0
    for i in range(n):  # strictly by index — never by OCR completion order
        text, kind, confs = native[i], PageKind.NATIVE, None
        page_notes: tuple[str, ...] = ()
        got = ocr_by_page.get(i)
        if got is not None:
            if got.failed:
                n_ocr_failed += 1
                page_notes = (f"{M_OCR_PAGE} to rasterize or read",)
            elif got.text.strip():
                text, kind, confs = got.text, PageKind.OCR, list(got.confs)
            else:
                kind = PageKind.OCR  # routed to OCR, recovered nothing
                confs = list(got.confs)
                n_ocr_blank += 1
            # The recovered tokens are appended to the TAIL, which is where the
            # Bates zone looks; ``BatesZone.tail_lines`` carries the matching
            # margin so the block can never evict a line the ordinary pass put
            # there. Their confidences join the page's, because they are text
            # on the page now and §4 Stage 2's threshold is measured over the
            # text the page actually carries.
            extra, extra_confs = reocr.get(i, ((), ()))
            if extra:
                have = {ln.strip() for ln in text.split("\n")}
                keep = [t for t in extra if t not in have]
                if keep:
                    text = (text.rstrip("\n") + "\n" if text.strip() else "") \
                        + "\n".join(keep)
                    confs = (confs or []) + [c for t, c in zip(extra, extra_confs)
                                             if t in keep]
                    n_footer_recovered += 1
        if i == 0 and photo:
            # The block describes the whole file, so it rides on page 1. When
            # the page also yielded read text the page stays OCR/NATIVE and
            # keeps its confidences — PHOTO is for a page whose only content
            # IS the deterministic block.
            has_read_text = bool(text.strip())
            text = (photo + "\n" + text) if has_read_text else photo
            if not has_read_text:
                kind, confs = PageKind.PHOTO, None
                # The relabelling loses the fact that this page WAS routed to
                # OCR and recovered nothing — ``make_page`` only adds that note
                # for a page handed to it as OCR. Disclosed here instead, so the
                # record still says what happened and :func:`ocr_yield` can
                # count the attempt. Without it, a corpus of photo-only PDFs run
                # against a dead engine reports zero attempts and no alarm.
                if got is not None and not got.failed:
                    page_notes = page_notes + (M_OCR_BLANK,)
        pages.append(make_page(i + 1, text, kind, confidences=confs,
                               conf_threshold=opt.conf_threshold, notes=page_notes))
    if n_ocr_failed:
        notes.append(f"{n_ocr_failed} page(s) could not be OCR'd; kept as empty pages")
    if n_ocr_blank:
        notes.append(f"{n_ocr_blank} page(s) routed to OCR recovered no text "
                     "(blank page, or nothing the engine could read)")
    if n_footer_recovered:
        # Disclosed, like every other bound in this module: an operator can see
        # how much of the production's numbering came from the second pass
        # rather than from the page's ordinary reading.
        notes.append(f"{n_footer_recovered} page(s) had a stamp-shaped token "
                     f"recovered by the targeted footer re-OCR "
                     f"({FOOTER_REOCR_DPI} dpi, "
                     f"{FOOTER_REOCR_BAND_PT / 72:.2f}in band, at most "
                     f"{FOOTER_BLOCK_MAX_LINES} token(s) per page)")
    if n_footer_declined:
        # The calibration bound, disclosed the same way every other bound in
        # this module is: neither band read a stamp on the pages it probed, so
        # the rest of the document was not re-read. Those pages are misses, and
        # an operator can see that they were never looked at rather than looked
        # at and found wanting.
        notes.append(f"{n_footer_declined} page(s) were not re-read by the "
                     f"targeted footer re-OCR: neither the footer nor the "
                     f"header band produced a stamp on the first "
                     f"{_FOOTER_PROBE_PAGES} page(s) probed in this document")
    return pages, notes


def _extract_docx(raw: bytes, opt: ExtractOptions) -> tuple[list[PageRecord], list[str]]:
    try:
        import docx  # python-docx
    except ImportError as exc:  # pragma: no cover — declared
        raise ExtractionError("Word support requires 'python-docx'.") from exc
    try:
        document = docx.Document(io.BytesIO(raw))
        parts = [p.text for p in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
    except Exception as exc:
        raise ExtractionError(f"Could not read Word document: {exc}") from exc
    note = "DOCX carries no page boundaries; emitted as one synthetic page"
    return synthetic_pages(["\n".join(parts)], notes=(note,)), [note]


def _xlsx_cell(v) -> str:
    """One cell → text. Dates render ISO so the date extractor anchors them."""
    if v is None:
        return ""
    if isinstance(v, datetime.datetime):
        # Drop a midnight time component so a pure date reads as 'YYYY-MM-DD'.
        return (v.date().isoformat() if v.time() == datetime.time(0, 0)
                else v.isoformat(sep=" "))
    if isinstance(v, datetime.date):
        return v.isoformat()
    return str(v)


def _extract_xlsx(raw: bytes, opt: ExtractOptions) -> tuple[list[PageRecord], list[str]]:
    """Workbook → one synthetic page per worksheet, tab-delimited.

    ``data_only=True`` yields the last-computed cell VALUES (not formulas) and
    ``read_only=True`` streams rows so a large register stays bounded in RAM.
    The row cap is global across the workbook and its truncation is disclosed
    both in the page text and as a document note — a silent cap would be a
    Principle-1 violation dressed as a performance guard.
    """
    try:
        import openpyxl
    except ImportError as exc:  # pragma: no cover — declared
        raise ExtractionError("Excel support requires 'openpyxl'.") from exc
    try:
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    except Exception as exc:
        raise ExtractionError(f"Could not read Excel workbook: {exc}") from exc

    blocks: list[str] = []
    notes: list[str] = []
    rows_emitted = 0
    truncated = False
    try:
        for ws in wb.worksheets:
            parts = [f"[sheet: {ws.title}]"]
            for row in ws.iter_rows(values_only=True):
                if rows_emitted >= _XLSX_MAX_ROWS:
                    truncated = True
                    break
                cells = [_xlsx_cell(c) for c in row]
                if not any(cells):
                    continue
                parts.append("\t".join(cells))
                rows_emitted += 1
            if truncated:
                parts.append(f"[... workbook truncated at {_XLSX_MAX_ROWS} rows]")
            blocks.append("\n".join(parts))
            if truncated:
                break
    finally:
        try:
            wb.close()
        except Exception:
            pass
    if truncated:
        notes.append(f"workbook truncated at {_XLSX_MAX_ROWS} rows; "
                     "later sheets were not read")
    notes.append("XLSX has no page boundaries; one synthetic page per worksheet")
    return synthetic_pages(blocks, notes=(notes[-1],)), notes


def _extract_xls(raw: bytes, opt: ExtractOptions) -> tuple[list[PageRecord], list[str]]:
    """Legacy .xls via xlrd — one synthetic page per sheet, same shape as XLSX."""
    try:
        import xlrd
    except ImportError as exc:  # pragma: no cover — declared
        raise ExtractionError("Legacy .xls support requires 'xlrd'.") from exc
    try:
        book = xlrd.open_workbook(file_contents=raw)
    except Exception as exc:
        raise ExtractionError(f"Could not read legacy Excel workbook: {exc}") from exc
    blocks: list[str] = []
    rows_emitted = 0
    truncated = False
    notes: list[str] = []
    for sheet in book.sheets():
        parts = [f"[sheet: {sheet.name}]"]
        for r in range(sheet.nrows):
            if rows_emitted >= _XLSX_MAX_ROWS:
                truncated = True
                break
            cells = ["" if c is None else str(c) for c in sheet.row_values(r)]
            if not any(c.strip() for c in cells):
                continue
            parts.append("\t".join(cells))
            rows_emitted += 1
        if truncated:
            parts.append(f"[... workbook truncated at {_XLSX_MAX_ROWS} rows]")
        blocks.append("\n".join(parts))
        if truncated:
            notes.append(f"workbook truncated at {_XLSX_MAX_ROWS} rows")
            break
    notes.append("XLS has no page boundaries; one synthetic page per worksheet")
    return synthetic_pages(blocks, notes=(notes[-1],)), notes


def _extract_csv(raw: bytes, opt: ExtractOptions) -> tuple[list[PageRecord], list[str]]:
    """CSV → tab-delimited text with an optional header marker.

    Delimiter auto-detection tries ``csv.Sniffer`` on the first 8 KB; if that
    fails it probes comma / semicolon / tab and picks whichever yields the most
    columns on average (ties broken in that order). Non-UTF-8 bytes fall back
    to latin-1. An empty CSV yields an empty page rather than raising — a
    zero-byte register is a fact about the production, not a failure.
    """
    import csv

    for enc in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            text = raw.decode(enc)
            break
        except Exception:
            continue
    else:  # pragma: no cover — latin-1 decodes every byte sequence
        text = raw.decode("latin-1", errors="replace")

    notes = ["CSV has no page boundaries; emitted as one synthetic page"]
    if not text.strip():
        return synthetic_pages([""], notes=(notes[0],)), notes

    sample = text[:8192]
    delim = ","
    try:
        delim = csv.Sniffer().sniff(sample, delimiters=",;\t").delimiter
    except Exception:
        best_avg = 0.0
        for d in (",", ";", "\t"):
            try:
                rows = list(csv.reader(sample.splitlines()[:200], delimiter=d))
                if rows:
                    avg = sum(len(r) for r in rows) / len(rows)
                    if avg > best_avg:
                        best_avg, delim = avg, d
            except Exception:
                pass

    def _looks_like_header(row: list[str]) -> bool:
        if not row or not all(c.strip() for c in row):
            return False
        return not any(re.fullmatch(r"\d+(\.\d+)?", c.strip()) for c in row)

    parts: list[str] = []
    rows_emitted = 0
    truncated = False
    first = True
    try:
        for row in csv.reader(text.splitlines(), delimiter=delim):
            if rows_emitted >= _CSV_MAX_ROWS:
                truncated = True
                break
            if first:
                first = False
                if _looks_like_header(row):
                    parts.append("[header: " + " | ".join(c.strip() for c in row) + "]")
                    continue
            if not any(c.strip() for c in row):
                continue
            parts.append("\t".join(row))
            rows_emitted += 1
    except Exception as exc:
        raise ExtractionError(f"Could not read CSV: {exc}") from exc
    if truncated:
        parts.append(f"[... CSV truncated at {_CSV_MAX_ROWS} rows]")
        notes.append(f"CSV truncated at {_CSV_MAX_ROWS} rows")
    return synthetic_pages(["\n".join(parts)], notes=(notes[0],)), notes


def _extract_pptx(raw: bytes, opt: ExtractOptions) -> tuple[list[PageRecord], list[str]]:
    """PowerPoint → one synthetic page per slide.

    Unlike the MIP 3.9 original, an empty slide still produces a page: slide 7
    of the source is slide 7 of the output, and Principle 1 accounts for it.
    """
    try:
        from pptx import Presentation  # python-pptx
    except ImportError as exc:  # pragma: no cover — declared
        raise ExtractionError("Presentation support requires 'python-pptx'.") from exc
    try:
        prs = Presentation(io.BytesIO(raw))
    except Exception as exc:
        raise ExtractionError(f"Could not read PowerPoint file: {exc}") from exc

    blocks: list[str] = []
    notes_failures: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append("\t".join(cells))
        # has_notes_slide FIRST: python-pptx's notes_slide property *creates* a
        # notes slide when none exists. Reading it unguarded manufactures parts
        # the source file never had, on every slide of every deck —
        # unacceptable in a tool whose claim is that every output is
        # mechanically derived from the source.
        try:
            if slide.has_notes_slide:
                for ph in slide.notes_slide.placeholders:
                    if ph.has_text_frame:
                        note_text = ph.text_frame.text.strip()
                        if note_text:
                            parts.append(f"[notes] {note_text}")
        except Exception as exc:
            # Was a bare ``pass``: a deck whose speaker notes would not read
            # lost them with no record anywhere, so the emitted text silently
            # differed from the source. Counted per slide and disclosed once.
            notes_failures.append(f"slide {len(blocks) + 1}: {exc}"[:120])
        blocks.append("\n".join(parts))
    note = "PPTX slides are emitted as synthetic pages, one per slide"
    out_notes = [note]
    if notes_failures:
        out_notes.append(
            f"{M_SLIDE_NOTES} on {len(notes_failures)} slide(s); their speaker "
            f"notes are absent from this document ({notes_failures[0]})")
    return synthetic_pages(blocks, notes=(note,)), out_notes


_RE_HTML_TAG = re.compile(r"<[^>]+>")


def _strip_html(s: str) -> str:
    import html as _html

    return _html.unescape(_RE_HTML_TAG.sub(" ", s or "")).strip()


def _extract_eml(raw: bytes, opt: ExtractOptions) -> tuple[list[PageRecord], list[str]]:
    """RFC-822 email → headers (From/To/Cc/Subject/Date) + body text.

    The ``Date`` header gets an extra ISO token appended so the date extractor
    anchors the message's own date; the body prefers text/plain, falling back
    to stripped HTML.
    """
    import email
    from email import policy
    from email.utils import parsedate_to_datetime

    notes: list[str] = []
    try:
        msg = email.message_from_bytes(raw, policy=policy.default)
    except Exception as exc:
        # Marked FINAL, not transient: the same bytes through the same parser
        # reach the same wall, so a serial re-read cannot improve it. Nothing
        # is deleted — the raw decode below carries every byte the file holds —
        # but the structure the rest of the pipeline reads (headers, the Date
        # anchor, the attachment list) is gone, and that is a gap the run must
        # be able to find mechanically rather than by reading prose.
        note = clip_message(f"{M_EML_PARSE}: decoded as raw text instead "
                            f"({exc}); headers, the Date anchor and the "
                            "attachment list were not recovered", 300)
        return (synthetic_pages([raw.decode("utf-8", errors="replace")], notes=(note,)),
                [note])
    parts: list[str] = []
    for h in ("From", "To", "Cc", "Subject"):
        v = msg.get(h)
        if v:
            parts.append(f"{h}: {v}")
    hdr_date = msg.get("Date")
    if hdr_date:
        try:
            iso = parsedate_to_datetime(hdr_date).date().isoformat()
        except Exception as exc:
            # The date sibling of B-3. The ISO token is the ONLY thing
            # ``dating.detect_dates`` can anchor an email's own date on —
            # RFC-2822 ("Tue, 16 Jul 2024 09:12:00 +0000") is not one of the
            # patterns it reads — so a Date header that will not parse used to
            # cost the message its date with no record at all.
            #
            # Deliberately NOT marked. Nothing DocIQ read is missing: the raw
            # header is still emitted verbatim below, and what is absent is a
            # derived convenience anchor. Marking it would put a malformed
            # sender's header into the run's evidence-loss tally, which is a
            # different and false claim. Disclosed, not marked.
            iso = ""
            notes.append(clip_message(
                f"the message's Date header {hdr_date!r} could not be parsed "
                f"({exc}); it is emitted verbatim but this document is not "
                "date-anchored on it", 300))
        parts.append(f"Date: {hdr_date}" + (f" ({iso})" if iso else ""))
    body = ""
    try:
        bp = msg.get_body(preferencelist=("plain", "html"))
        if bp is not None:
            body = bp.get_content()
            if bp.get_content_subtype() == "html":
                body = _strip_html(body)
    except Exception as exc:
        # Codex review #1, B-3. This was ``body = ""`` under a bare
        # ``except``: a supported email whose body would not decode came back
        # with its headers, no body, no note, no marker and a FULL status. The
        # message text — the whole evidentiary point of an email — was deleted
        # and the run reported success, so the walker's serial-retry registry
        # never saw the file either.
        body = ""
        notes.append(clip_message(
            f"{M_EML_BODY}: the message body could not be decoded ({exc}); "
            "this document carries its headers only", 300))
    if body and body.strip():
        parts.append("")
        parts.append(body.strip())
    page_note = "email carries no page boundaries; emitted as one synthetic page"
    return (synthetic_pages(["\n".join(parts)], notes=(page_note,)),
            [page_note] + notes)


def _extract_msg(raw: bytes, opt: ExtractOptions) -> tuple[list[PageRecord], list[str]]:
    """Outlook ``.msg`` → headers + body via ``extract-msg``.

    The library needs a real path, so a scratch file is unavoidable. It goes
    under the caller's working folder when one was supplied (§10) and is
    unlinked in a ``finally`` either way.
    """
    try:
        import extract_msg
    except ImportError as exc:  # pragma: no cover — declared
        raise ExtractionError("Outlook .msg support requires 'extract-msg'.") from exc
    import tempfile

    scratch = opt.scratch_dir
    if scratch is not None:
        scratch.mkdir(parents=True, exist_ok=True)
    tmp = None
    try:
        with tempfile.NamedTemporaryFile(
            suffix=".msg", delete=False,
            dir=str(scratch) if scratch is not None else None,
        ) as tf:
            tf.write(raw)
            tmp = tf.name
        m = extract_msg.Message(tmp)
        parts: list[str] = []
        for label, val in (("From", getattr(m, "sender", None)),
                           ("To", getattr(m, "to", None)),
                           ("Cc", getattr(m, "cc", None)),
                           ("Subject", getattr(m, "subject", None)),
                           ("Date", getattr(m, "date", None))):
            if val:
                parts.append(f"{label}: {val}")
        body = getattr(m, "body", None)
        if body:
            parts.append("")
            parts.append(str(body).strip())
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"Could not read Outlook .msg: {exc}") from exc
    finally:
        if tmp:
            try:
                os.unlink(tmp)
            except Exception:
                pass
    note = "Outlook message carries no page boundaries; one synthetic page"
    return synthetic_pages(["\n".join(parts)], notes=(note,)), [note]


def _extract_text(raw: bytes, opt: ExtractOptions) -> tuple[list[PageRecord], list[str]]:
    note = "plain text carries no page boundaries; emitted as one synthetic page"
    return (synthetic_pages([raw.decode("utf-8", errors="replace")], notes=(note,)),
            [note])


# ---------------------------------------------------------------------------
# ZIP — members become child documents, not concatenated text
# ---------------------------------------------------------------------------


@dataclass(frozen=True, slots=True)
class ZipMember:
    name: str
    raw: bytes
    order: int
    """Position in ``infolist()`` order. Deterministic child ID assignment
    (D-04) depends on it, so it is carried, not recomputed."""


@dataclass(frozen=True, slots=True)
class ZipExpansion:
    members: tuple[ZipMember, ...] = ()
    notes: tuple[str, ...] = ()


def expand_eml_attachments(raw: bytes) -> ZipExpansion:
    """Attachments of an RFC-822 email, as child members.

    §3 requires MSG/EML attachments to be "extracted as child documents linked
    to the parent message ID" — Tier 1, not optional. ``_extract_eml`` only
    ever produced the message's own headers+body page; nothing walked
    ``iter_attachments()``, so every attachment on every email in a matter
    vanished with no record, no note, and no line in the Unsupported list —
    a silent deletion Principle 1 forbids outright. This is the missing half.

    A zip attachment is flattened one level via :func:`expand_zip`, the same
    treatment a zip-inside-a-zip already gets, so "attach the production as a
    zip" does not reopen the hole this closes.
    """
    import email
    from email import policy

    try:
        msg = email.message_from_bytes(raw, policy=policy.default)
    except Exception as exc:
        # Codex review #1, B-3. This was a bare ``return ZipExpansion()``: the
        # parent message was emitted with ZERO attachments, no note and no
        # marker, so an email carrying the production itself looked like an
        # email carrying nothing. M_ATTACH_ENUM is the transient marker the
        # walker's retry registry keys on, which is what makes the failure
        # re-read rather than written off.
        return ZipExpansion((), (clip_message(
            f"{M_ATTACH_ENUM}: the message envelope would not parse ({exc}), "
            "so NO attachment of this email was brought in", 200),))

    raw_members: list[tuple[str, bytes]] = []
    notes: list[str] = []
    try:
        parts = list(msg.iter_attachments())
    except Exception as exc:
        return ZipExpansion((), (f"{M_ATTACH_ENUM}: {exc}"[:200],))
    for part in parts:
        try:
            name = part.get_filename() or f"attachment_{len(raw_members) + 1}"
            payload = part.get_payload(decode=True)
        except Exception as exc:
            notes.append(f"{M_ATTACH_READ}: {exc}"[:200])
            continue
        if payload is None:
            # Disclosed before, but with no marker of any kind — so nothing
            # downstream could find it mechanically and it sat outside both the
            # retry registry and the accounting tally (Codex review #1, B-3).
            # FINAL rather than transient: a part with no decodable payload
            # decodes to nothing on the second attempt too.
            notes.append(f"{M_ATTACH_SKIPPED}: attachment '{name}' had no "
                         "decodable payload; it is named here and its bytes "
                         "are not in the corpus")
            continue
        raw_members.append((name, payload))

    members: list[ZipMember] = []
    for name, payload in raw_members:
        if _ext(name) == ".zip":
            try:
                inner = expand_zip(payload)
            except Exception as exc:
                notes.append(f"{M_ZIP_ATTACH}: attachment '{name}' is a zip "
                             f"that could not be read: {exc}"[:200])
                continue
            notes.extend(f"{name}: {n}" for n in inner.notes)
            for m in inner.members:
                members.append(ZipMember(f"{name}/{m.name}", m.raw, len(members)))
        else:
            members.append(ZipMember(name, payload, len(members)))
    return ZipExpansion(tuple(members), tuple(notes))


def expand_msg_attachments(raw: bytes, scratch_dir: Path | None) -> ZipExpansion:
    """Attachments of an Outlook ``.msg``, as child members. See
    :func:`expand_eml_attachments` — same requirement, same prior gap.

    ``extract-msg`` needs a real path, same as ``_extract_msg``; the scratch
    file goes under the caller's working folder (§10) and is unlinked either
    way. An embedded-message attachment (an ``.msg`` inside a ``.msg``, which
    the library returns as a nested ``Message`` rather than bytes) is
    disclosed rather than silently skipped: it is real content, just not one
    this pass can flatten without recursing into a second temp-file dance.
    """
    try:
        import extract_msg
    except ImportError as exc:  # pragma: no cover — declared
        raise ExtractionError("Outlook .msg support requires 'extract-msg'.") from exc
    import tempfile

    if scratch_dir is not None:
        scratch_dir.mkdir(parents=True, exist_ok=True)
    tmp = None
    notes: list[str] = []
    raw_members: list[tuple[str, bytes]] = []
    try:
        with tempfile.NamedTemporaryFile(
            suffix=".msg", delete=False,
            dir=str(scratch_dir) if scratch_dir is not None else None,
        ) as tf:
            tf.write(raw)
            tmp = tf.name
        m = extract_msg.Message(tmp)
        for i, att in enumerate(getattr(m, "attachments", None) or []):
            try:
                name = (getattr(att, "longFilename", None)
                        or getattr(att, "shortFilename", None)
                        or f"attachment_{i + 1}")
                data = getattr(att, "data", None)
            except Exception as exc:
                notes.append(f"{M_ATTACH_READ}: {exc}"[:200])
                continue
            if isinstance(data, (bytes, bytearray)):
                raw_members.append((name, bytes(data)))
            else:
                # An embedded .msg (Outlook nests a Message object, not
                # bytes) or an unreadable attachment kind. Disclosed, not
                # dropped: the operator sees that content exists and was not
                # brought in, rather than the run looking complete.
                notes.append(f"{M_ATTACH_SKIPPED}: attachment '{name}' is an "
                             "embedded message or an unsupported attachment "
                             "kind; it is named here and its bytes are not in "
                             "the corpus")
    except Exception as exc:
        return ZipExpansion((), (f"{M_MSG_ATTACH}: {exc}"[:200],))
    finally:
        if tmp:
            try:
                os.unlink(tmp)
            except Exception:
                pass

    members: list[ZipMember] = []
    for name, payload in raw_members:
        if _ext(name) == ".zip":
            try:
                inner = expand_zip(payload)
            except Exception as exc:
                notes.append(f"{M_ZIP_ATTACH}: attachment '{name}' is a zip "
                             f"that could not be read: {exc}"[:200])
                continue
            notes.extend(f"{name}: {n}" for n in inner.notes)
            for m2 in inner.members:
                members.append(ZipMember(f"{name}/{m2.name}", m2.raw, len(members)))
        else:
            members.append(ZipMember(name, payload, len(members)))
    return ZipExpansion(tuple(members), tuple(notes))


def expand_zip(raw: bytes, depth: int = 0) -> ZipExpansion:
    """Expand a (possibly nested) ZIP into flat member byte blobs.

    Members are read into memory only — never written to disk, so archive
    path-traversal is moot — and total uncompressed bytes, member count and
    nesting depth are all capped against a malicious or accidental bomb. Every
    cap that bites is disclosed as a note; §"no silent caps" is not satisfied
    by a guard that quietly stops early.

    Unlike MIP 3.9, members are returned rather than concatenated: the contract
    models an archive member as its own :class:`DocumentRecord` with a
    ``parent_doc_id`` and a ``container_order``.
    """
    import zipfile

    notes: list[str] = []
    if depth > _ZIP_MAX_DEPTH:
        return ZipExpansion((), (f"zip nesting deeper than {_ZIP_MAX_DEPTH} "
                                 "levels was not expanded",))
    try:
        zf = zipfile.ZipFile(io.BytesIO(raw))
    except Exception as exc:
        raise ExtractionError(f"Could not read ZIP archive: {exc}") from exc

    members: list[ZipMember] = []
    total = 0
    cap_bytes = _ZIP_MAX_MB * 1024 * 1024
    with zf:
        for info in zf.infolist():
            if info.is_dir():
                continue
            if len(members) >= _ZIP_MAX_MEMBERS:
                notes.append(f"archive truncated at {_ZIP_MAX_MEMBERS} members; "
                             "later members were not read")
                break
            if info.file_size and total + info.file_size > cap_bytes:
                notes.append(f"archive truncated at {_ZIP_MAX_MB} MB uncompressed; "
                             f"'{info.filename}' and later members were not read")
                break
            try:
                blob = zf.read(info)
            except Exception as exc:
                notes.append(f"{M_ZIP_MEMBER}: '{info.filename}': "
                             f"{str(exc)[:120]}")
                continue
            total += len(blob)
            if _ext(info.filename) == ".zip":
                inner = expand_zip(blob, depth + 1)
                notes.extend(f"{info.filename}: {n}" for n in inner.notes)
                for m in inner.members:
                    members.append(ZipMember(f"{info.filename}/{m.name}", m.raw,
                                             len(members)))
                continue
            members.append(ZipMember(info.filename, blob, len(members)))
    return ZipExpansion(tuple(members), tuple(notes))


# ---------------------------------------------------------------------------
# Content-sniffing recovery (extension-mismatch)
# ---------------------------------------------------------------------------

_MAGIC_PDF = b"%PDF"
_MAGIC_ZIP = b"PK\x03\x04"
_MAGIC_OLE = b"\xd0\xcf\x11\xe0"

# The zip signature is shared by every OOXML format, so the whole family is
# probed, plain ZIP last.
_SNIFF_CHAINS = {
    "pdf": [(".pdf", "PDF")],
    "zip": [(".docx", "Word"), (".xlsx", "Excel"), (".pptx", "PowerPoint")],
    "ole": [(".msg", "Outlook .msg"), (".xls", "legacy Excel")],
}
_SNIFF_LABELS = {"pdf": "PDF", "zip": "a zip-family container",
                 "ole": "a legacy OLE container"}
# Extensions sharing one extractor — a failed .xlsm must not retry .xlsx.
_EXT_ALIASES = {".xlsm": ".xlsx", ".email": ".eml", ".md": ".txt", ".log": ".txt"}


def _ext(filename: str) -> str:
    name = (filename or "").lower()
    dot = name.rfind(".")
    return name[dot:] if dot >= 0 else ""


# Windows drive paths, UNC paths, and POSIX absolute paths.
_ABS_PATH = re.compile(r"(?:[A-Za-z]:[\\/]|\\\\|(?<![\w.])/)[^\s\"'<>|]*")


def clip_message(msg: str, limit: int) -> str:
    """Truncate a message and SAY that it was truncated.

    Every error string in the pipeline is length-bounded so one pathological
    parser cannot write a megabyte into the log. A bound that removes text
    without a mark is a silent cap, which is the thing the standing rule
    forbids — so the mark is not decoration.
    """
    if msg is None or len(msg) <= limit:
        return msg
    return msg[:limit].rstrip() + f" […truncated at {limit} chars]"


def sanitize_message(msg: str) -> str:
    """Strip absolute paths out of a message that will reach a record.

    An error string is hashed content: it lands in ``DocumentRecord.error``,
    which lands in the log's ``content`` section. A parser that reports
    ``C:\\Users\\...\\Temp\\tmp61yhcl7p\\x.msg`` therefore puts a per-run
    random string inside the byte-identical claim. Reducing every absolute
    path to its final component fixes the whole class at the one place every
    message passes, rather than auditing each ``f"...{exc}"`` for ever.
    """
    if not msg:
        return msg

    def _basename(m: re.Match[str]) -> str:
        return m.group(0).replace("\\", "/").rstrip("/").rsplit("/", 1)[-1] or "<path>"

    return _ABS_PATH.sub(_basename, msg)


def sniff_kind(raw: bytes) -> str:
    """Magic-byte kind (``pdf`` / ``zip`` / ``ole``), or ``""``."""
    if raw.startswith(_MAGIC_PDF):
        return "pdf"
    if raw.startswith(_MAGIC_ZIP):
        return "zip"
    if raw.startswith(_MAGIC_OLE):
        return "ole"
    return ""


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".xlsm": _extract_xlsx,
    ".xls": _extract_xls,
    ".csv": _extract_csv,
    ".pptx": _extract_pptx,
    ".eml": _extract_eml,
    ".email": _extract_eml,
    ".msg": _extract_msg,
    ".txt": _extract_text,
    ".md": _extract_text,
    ".log": _extract_text,
}


def _dispatch(ext: str, raw: bytes, opt: ExtractOptions):
    fn = _EXTRACTORS.get(ext)
    if fn is None:
        raise ExtractionError(f"No Tier-1 extractor for '{ext}'.")
    return fn(raw, opt)


def _retry_by_content(ext: str, raw: bytes, opt: ExtractOptions, original: Exception):
    """Retry with the content-sniffed extractor after ``ext``'s own raised.

    Litigation productions routinely deliver files under the wrong extension —
    PDF bytes named .docx, Word files named .pdf — and an extension-only
    dispatch skips them with "File is not a zip file". Recovery is annotated so
    the mismatch reaches the audit trail rather than being quietly fixed.
    """
    kind = sniff_kind(raw)
    if not kind:
        raise original
    canon = _EXT_ALIASES.get(ext, ext)
    tried: list[str] = []
    for retry_ext, label in _SNIFF_CHAINS[kind]:
        if retry_ext == canon:
            continue  # the extension-selected extractor already failed
        tried.append(label)
        try:
            pages, notes = _dispatch(retry_ext, raw, opt)
        except Exception:
            continue
        return pages, notes + [f"extension {ext} but content is "
                               f"{_SNIFF_LABELS[kind]}; recovered via "
                               f"{label} extractor"]
    raise ExtractionError(
        f"{original} (content sniffed as {_SNIFF_LABELS[kind]}; retry via "
        + (", ".join(tried) if tried else "no alternate extractor")
        + " also failed)"
    ) from original


# ---------------------------------------------------------------------------
# Public entry points
# ---------------------------------------------------------------------------


def is_tier1(ext: str) -> bool:
    return ext.lower() in TIER1_EXTENSIONS


def tier2_hint(ext: str) -> str:
    """The remediation hint for a Tier-2 or unknown extension (§3, D-02)."""
    return TIER2_EXTENSIONS.get(ext.lower(), UNKNOWN_HINT)


def extract(filename: str, raw: bytes,
            opt: ExtractOptions | None = None) -> ExtractedDoc:
    """Extract one document's pages. Never raises.

    A failure is data — an :class:`ExtractedDoc` with FAILED status and an
    actionable message — because a single unreadable file in a 9,000-file
    production must not abort the run.
    """
    opt = opt or ExtractOptions()
    ext = _ext(filename)
    if ext == ".zip":
        raise ExtractionError(
            "ZIP is expanded by the walker into child documents; "
            "call expand_zip() instead of extract().")
    if not is_tier1(ext):
        return ExtractedDoc(status=ProcessingStatus.UNSUPPORTED,
                            error=tier2_hint(ext))
    if not raw:
        return ExtractedDoc(status=ProcessingStatus.FAILED, error="empty file")
    try:
        pages, notes = _dispatch(ext, raw, opt)
    except Exception as exc:
        try:
            pages, notes = _retry_by_content(ext, raw, opt, exc)
        except Exception as exc2:
            return ExtractedDoc(status=ProcessingStatus.FAILED,
                                error=clip_message(sanitize_message(str(exc2)), 400))
    notes = [sanitize_message(n) for n in notes]
    # §4 Stage 2: the flag is driven by the page's MEAN confidence against the
    # run threshold. A page that failed OCR outright flags too — it is exactly
    # the case a human must look at, and it would otherwise pass as FULL
    # because it has no confidence to be below anything. A page that OCR'd
    # cleanly to nothing does NOT flag: a blank page inside a native PDF is
    # ordinary, and flagging it would train the operator to ignore the flag.
    # It is still disclosed, as a page note and a document note.
    flagged = any(
        (p.ocr_conf is not None and p.ocr_conf < opt.conf_threshold)
        or any(n.startswith(M_OCR_PAGE) for n in p.notes)
        for p in pages
    )
    status = (ProcessingStatus.PARTIAL_OCR_FLAGGED if flagged
              else ProcessingStatus.FULL)
    return ExtractedDoc(pages=tuple(pages), notes=tuple(notes), status=status)

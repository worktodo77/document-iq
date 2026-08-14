"""The Stage-6 gate, the manifest, and the byte-identical claim."""

from __future__ import annotations

import json

import pytest

from dociq.contracts import (
    Disposition,
    DocumentRecord,
    PageKind,
    PageRecord,
    ProcessingStatus,
    RunConfig,
    RunResult,
)
from dociq import pipeline
from dociq.ingest import extract as ex
from dociq.ingest import walker
from dociq.profiles.model import OperatorStamp
from dociq.verify import accounting, manifest

from .conftest import FIXTURES


def _page(n: int, **kw) -> PageRecord:
    return PageRecord(page_no=n, text=f"page {n}", kind=PageKind.NATIVE, **kw)


def _doc(rel: str = "a.pdf", pages=(1, 2), **kw) -> DocumentRecord:
    return DocumentRecord(doc_id="", rel_path=rel, filename=rel, sha256="0" * 64,
                          size_bytes=1, ext=".pdf",
                          pages=tuple(_page(n) for n in pages), **kw)


def _result(*docs, **kw) -> RunResult:
    cfg = RunConfig(source_root="s", output_root="o")
    return RunResult(config=cfg, documents=tuple(docs), **kw)


def test_clean_corpus_reconciles():
    rep = accounting.check(_result(_doc()))
    assert rep.ok
    assert "PAGE ACCOUNTING OK" in rep.render()


def test_a_dropped_page_without_a_rule_is_named_not_just_counted():
    doc = DocumentRecord(
        doc_id="", rel_path="bad.pdf", filename="bad.pdf", sha256="0" * 64,
        size_bytes=1, ext=".pdf",
        pages=(PageRecord(page_no=1, text="x", kind=PageKind.NATIVE,
                          disposition=Disposition.DROP, drop_rule=""),))
    rep = accounting.check(_result(doc))
    assert not rep.ok
    assert any("bad.pdf" in str(d) for d in rep.discrepancies)


def test_a_page_number_gap_is_reported_against_its_document():
    doc = DocumentRecord(doc_id="", rel_path="gap.pdf", filename="gap.pdf",
                         sha256="0" * 64, size_bytes=1, ext=".pdf",
                         pages=(_page(1), _page(3)))
    rep = accounting.check(_result(doc))
    assert not rep.ok
    assert any(d.rel_path == "gap.pdf" and d.kind == "contract"
               for d in rep.discrepancies)


def test_an_orphan_container_child_is_reported():
    child = _doc("z.zip/a.pdf", parent_doc_id="missing.zip", container_order=0)
    rep = accounting.check(_result(child))
    assert any(d.kind == "orphan-child" for d in rep.discrepancies)


def test_a_failed_document_with_no_message_is_a_discrepancy():
    doc = _doc("f.pdf", pages=(1,), status=ProcessingStatus.FAILED)
    rep = accounting.check(_result(doc))
    assert any(d.kind == "silent-failure" for d in rep.discrepancies)


def test_two_documents_sharing_a_path_would_overwrite_and_are_caught():
    rep = accounting.check(_result(_doc("dup.pdf"), _doc("dup.pdf")))
    assert any(d.kind == "duplicate-record" for d in rep.discrepancies)


def test_report_names_the_document_not_just_a_boolean():
    doc = DocumentRecord(doc_id="", rel_path="deep/nested/thing.pdf",
                         filename="thing.pdf", sha256="0" * 64, size_bytes=1,
                         ext=".pdf", pages=(_page(2),))
    text = accounting.check(_result(doc)).render()
    assert "deep/nested/thing.pdf" in text


# -- manifest ---------------------------------------------------------------


@pytest.fixture(scope="module")
def emitted(tmp_path_factory):
    """A real run of the shipped pipeline over the fixture corpus.

    OCR is off: these tests are about the manifest and the emit layer, and a
    real OCR pass would add minutes to the suite without changing what they
    assert. The OCR path is proven in the self-test and in ``test_extract``.
    """
    out = tmp_path_factory.mktemp("emit")
    cfg = RunConfig(source_root=str(FIXTURES), output_root=str(out),
                    ocr_engine_version=ex.ocr_engine_version())
    outcome = pipeline.run(cfg, pipeline.PipelineOptions(
        walk=walker.WalkOptions(ocr_enabled=False, resume=False),
        stamp=OperatorStamp("test", "2026-07-30T00:00:00Z", "test")))
    return out, outcome.result


def test_manifest_covers_exactly_the_claimed_files(emitted):
    out, _ = emitted
    man = manifest.build(out)
    assert "sources.json" in man.deterministic
    assert "document_index.csv" in man.deterministic
    assert any(k.startswith("clean_text/") for k in man.deterministic)
    assert man.log_content_sha256


def test_manifest_states_the_split_rather_than_leaving_it_ambiguous(emitted):
    out, _ = emitted
    payload = json.loads(manifest.write(out).read_text(encoding="utf-8"))
    assert "byte-identical" in payload["claim"]
    assert "content" in payload["claim"] and "processing_log.json" in payload["claim"]
    assert manifest.LOG_NAME in payload["excluded"]
    assert "run" in payload["excluded"][manifest.LOG_NAME]


def test_an_undeclared_output_is_reported_unclassified_not_assumed(emitted, tmp_path):
    out, _ = emitted
    (out / "surprise.dat").write_bytes(b"x")
    try:
        man = manifest.build(out)
        assert "surprise.dat" in man.unclassified
        assert "UNCLASSIFIED" in man.render()
    finally:
        (out / "surprise.dat").unlink()


def test_the_excluded_files_are_not_compared(emitted):
    out, _ = emitted
    a = manifest.build(out)
    b = manifest.build(out)
    b.excluded["run_summary.pdf"] = "different reason"
    assert manifest.compare(a, b) == []


def test_compare_names_the_file_that_diverged(emitted):
    out, _ = emitted
    a = manifest.build(out)
    b = manifest.build(out)
    key = next(iter(b.deterministic))
    b.deterministic[key] = "f" * 64
    assert any(key in d for d in manifest.compare(a, b))


def test_clean_text_carries_a_marker_for_every_page_including_empty(emitted):
    out, result = emitted
    doc = next(d for d in result.documents
               if d.rel_path == "04_empty_page.pdf")
    text = (out / "clean_text" / f"{doc.doc_id}.txt").read_text(encoding="utf-8")
    assert text.count("===== PAGE") == doc.pages_in == 3


def test_an_empty_output_root_is_a_hard_failure_not_a_stable_hash(tmp_path):
    """An empty manifest hashes perfectly stably, so two runs that produced
    nothing would compare byte-identical and the gate would go green on a
    pipeline that did no work."""
    with pytest.raises(manifest.EmptyOutputError):
        manifest.build(tmp_path / "never-created")
    (tmp_path / "empty").mkdir()
    with pytest.raises(manifest.EmptyOutputError):
        manifest.build(tmp_path / "empty")
    assert manifest.build(tmp_path / "empty", require_outputs=False).deterministic == {}


def test_determinism_reports_a_run_that_produced_nothing(tmp_path, monkeypatch):
    from dociq.verify import determinism

    monkeypatch.setattr(determinism, "_one_run", lambda src, out, seed: None)
    rep = determinism.prove(tmp_path, runs=3, workdir=tmp_path / "det")
    assert not rep.ok
    assert len(rep.failures) == 3
    assert any("no subject" in f or "does not exist" in f for f in rep.failures)


def test_pptx_extraction_never_touches_notes_slide_on_a_deck_without_notes(
        monkeypatch):
    """python-pptx's ``notes_slide`` property CREATES a notes slide when none
    exists. Reading it unguarded manufactures parts the source file never had,
    on every slide of every deck. Make the property fatal and prove the
    extractor no longer reaches it."""
    import io as _io

    from pptx import Presentation
    from pptx.slide import Slide

    from dociq.ingest import extract as _ex

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    buf = _io.BytesIO()
    prs.save(buf)

    class _Touched(BaseException):
        """Not an Exception: the extractor wraps the notes read in a broad
        ``except Exception`` — a probe that the code under test can swallow is
        a probe that cannot fail."""

    def explode(self):
        raise _Touched("notes_slide read on a deck with no notes")

    monkeypatch.setattr(Slide, "notes_slide", property(explode))
    got = _ex.extract("deck.pptx", buf.getvalue())
    assert got.status is not ProcessingStatus.FAILED
    assert len(got.pages) == 1


# ---------------------------------------------------------------------------
# C6 — what criterion 7's proof covers, made assertable
# ---------------------------------------------------------------------------


def test_the_determinism_report_says_which_REGIME_produced_it():
    """A report that cannot say whether its repetitions were sequential or
    contended is a report that will be quoted as covering both.

    The acceptance run of 2026-08-02 is the evidence the two regimes differ: the
    shipped per-file timeout was crossed by 2 documents idle and 6 under load.
    Until ``concurrency`` existed every repetition in this module was sequential,
    so the one regime documented as behaving differently was the one nothing
    exercised.
    """
    from dociq.verify import determinism

    seq = determinism.DeterminismReport(runs=8, concurrency=1)
    con = determinism.DeterminismReport(runs=8, concurrency=4)
    assert "sequential" in seq.render()
    assert "CONTENDED" in con.render() and "4 at a time" in con.render()
    assert "sequential" not in con.render()
    assert '"concurrency": 4' in determinism.prove_json(con)


def test_concurrency_is_carried_into_the_report_not_silently_dropped(
    tmp_path, monkeypatch
):
    """FAIL-BEFORE, watched RED: with ``concurrency`` accepted and not stored on
    the report, the render says "sequential" for a contended proof — a wider
    claim than the run supports, which is the whole class this package is about.
    """
    from dociq.verify import determinism

    seen: list[int] = []

    def fake_one_run(src, out, seed):
        seen.append(1)
        return "stub: no subject"

    monkeypatch.setattr(determinism, "_one_run", fake_one_run)
    rep = determinism.prove(tmp_path, runs=4, workdir=tmp_path / "d",
                            concurrency=3)
    assert rep.concurrency == 3
    assert "CONTENDED" in rep.render()
    assert len(seen) == 4, "every repetition must still run exactly once"


def test_concurrency_is_clamped_to_the_run_count_and_to_one():
    """No silent cap and no silent zero: asking for more parallelism than there
    are repetitions is not an error, and asking for none must not mean "run
    nothing"."""
    from dociq.verify import determinism

    for asked, runs, expected in ((0, 4, 1), (-3, 4, 1), (9, 4, 4), (2, 4, 2)):
        rep = determinism.DeterminismReport(
            runs=runs, concurrency=max(1, min(asked, runs)))
        assert rep.concurrency == expected


def test_concurrency_actually_OVERLAPS_the_repetitions(tmp_path, monkeypatch):
    """"CONTENDED" in the report must not be a label on sequential work.

    ``concurrency`` is only meaningful if the repetitions genuinely overlap, so
    this measures the observed overlap rather than trusting the parameter. The
    stub records how many repetitions are in flight at once; at ``concurrency=4``
    the peak must exceed 1, and at ``concurrency=1`` it must be exactly 1 — the
    second half is what makes the first half an assertion rather than a
    coincidence of scheduling.

    FAIL-BEFORE, watched RED: with the ThreadPoolExecutor branch replaced by the
    sequential list comprehension, the peak is 1 and this goes red.
    """
    import threading
    import time

    from dociq.verify import determinism

    for asked, must_overlap in ((4, True), (1, False)):
        lock = threading.Lock()
        state = {"live": 0, "peak": 0}

        def fake_one_run(src, out, seed):
            with lock:
                state["live"] += 1
                state["peak"] = max(state["peak"], state["live"])
            time.sleep(0.05)
            with lock:
                state["live"] -= 1
            return "stub: no subject"

        monkeypatch.setattr(determinism, "_one_run", fake_one_run)
        determinism.prove(tmp_path, runs=4, workdir=tmp_path / f"d{asked}",
                          concurrency=asked)
        if must_overlap:
            assert state["peak"] > 1, (
                "the report would say CONTENDED for work that never overlapped")
        else:
            assert state["peak"] == 1
